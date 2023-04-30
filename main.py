import cv2
import numpy as np
import os
import sys
import tensorflow as tf
import pptx
from pptx.util import Inches
import csv
from PIL import Image

from sklearn.model_selection import train_test_split

EPOCHS = 10
IMG_WIDTH = 30
IMG_HEIGHT = 30
NUM_CATEGORIES = 43
TEST_SIZE = 0.4

data_directory = "gtsrb"
model_name = "model"
csvDir = "output.csv"
id_dir = "id_dir"

def textToBool(inp):
    return inp in ["True", "true"]

def main():
    # Check command-line arguments
    if len(sys.argv) not in [2, 3]:
        sys.exit("Usage: python main.py (bool)train_model")

    if (textToBool(sys.argv[1])):
        images, labels = load_data(data_directory)

        # Split data into training and testing sets
        labels = tf.keras.utils.to_categorical(labels)
        x_train, x_test, y_train, y_test = train_test_split(
            np.array(images), np.array(labels), test_size=TEST_SIZE
        )

        # Get a compiled neural network
        model = get_model()

        # Fit model on training data
        model.fit(x_train, y_train, epochs=EPOCHS)

        # Evaluate neural network performance
        model.evaluate(x_test,  y_test, verbose=2)

        model.save(model_name)
        print(f"Model saved to {model_name}.")
    else:
        model = tf.keras.models.load_model(model_name)
        imgs = os.listdir(id_dir)
        prs = pptx.Presentation()
        lyt=prs.slide_layouts[0] # choosing a slide layout

        for image in imgs:
            if image == ".DS_Store":
                continue

            imgPath = os.path.join(id_dir, image)

            if imgPath[:-4] != ".ppm":
                im = Image.open(imgPath)
                ppmPath = f"{imgPath}.ppm"
                im.convert("RGB").save(ppmPath, "PPM")
                os.rename(imgPath, os.path.join("wrong_formats", image))
                imgPath = ppmPath

            img = cv2.imread(imgPath, cv2.IMREAD_COLOR)
            img = cv2.resize(img, (IMG_WIDTH, IMG_HEIGHT))
            op = model.predict(np.array([img]))

            newPath = os.path.join(id_dir, f"{str(imgs.index(image))}.jpg")
            im = Image.open(imgPath)
            im.convert("RGB").save(newPath, "JPEG")

            slide=prs.slides.add_slide(lyt) #New slide
            # title=slide.shapes.title
            imag=slide.shapes.add_picture(newPath, Inches(0), Inches(0)) #Image
            subtitle=slide.placeholders[1]
            # title.text=image
            subtitle.text= str(np.argmax(op)) #Image classification
            os.remove(newPath)
        prs.save("plants.pptx")

def load_data(data_dir):
    images = []
    labels = []
    gtsrb = os.listdir(data_dir)
    for folder in gtsrb:
      if not folder.startswith("."):
        path = os.path.join(data_dir, folder)
        for image in os.listdir(path):
          img = cv2.imread(os.path.join(path, image), cv2.IMREAD_COLOR)
          img = cv2.resize(img, (IMG_WIDTH, IMG_HEIGHT))
          images.append(img)
          labels.append(int(folder))
    return images, labels

def get_model():
    model = tf.keras.models.Sequential(
      [
        tf.keras.layers.Conv2D(32, (3, 3), activation="relu", input_shape=(IMG_WIDTH, IMG_HEIGHT, 3)),
        tf.keras.layers.MaxPooling2D(pool_size=(3, 3)),
        tf.keras.layers.Flatten(),
        #relu, sigmoid, softmax

        tf.keras.layers.Dropout(0.2),
        tf.keras.layers.Dense(NUM_CATEGORIES * 16, activation="relu"),
        tf.keras.layers.Dropout(0.2),
        tf.keras.layers.Dense(NUM_CATEGORIES * 8, activation="relu"),

        tf.keras.layers.Dense(NUM_CATEGORIES, activation="softmax")
      ]
    )

    model.compile(
      optimizer="adam",
      loss="binary_crossentropy",
      metrics=["accuracy"]
    )

    return model

if __name__ == "__main__":
    main()
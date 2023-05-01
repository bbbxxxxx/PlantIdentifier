import tensorflow as tf
import os
from PIL import Image
import json
import requests
import pptx
from pptx.util import Inches
import numpy as np
import csv
import sys
import cv2

from json import JSONEncoder
from selenium import webdriver
from webdriver_manager.chrome import ChromeDriverManager
from selenium.webdriver.common.by import By
from selenium.common.exceptions import NoSuchElementException

import tree_maker as tm
from bs4 import BeautifulSoup as BS
import re

valDir = "values.json"
plantsList = "plants.csv"
id_dir = "id_dir"
searchUrl = "https://www.biolib.cz/cz/main/"

resultDict = {}
resultTree = []

class plantEncoder(JSONEncoder):
        def default(self, o):
            return {o.__class__.__name__: o.__dict__}

class plant:
    def __init__(self, kingdom, phylum, plantClass, order, family, genus, czechName, latinName):
        self.kingdom = kingdom
        self.phylum = phylum
        self.plantClass = plantClass
        self.order = order
        self.family = family
        self.genus = genus
        self.czechName = czechName
        self.latinName = latinName

def produce_output():
    imgs = os.listdir(id_dir)
    prs = pptx.Presentation()
    lyt=prs.slide_layouts[0] # choosing a slide layout

    for image in imgs:
        if image in [".DS_Store", ".gitkeep"]:
            continue

        imgPath = os.path.join(id_dir, image)

        # if imgPath[:-4] != ".ppm":
        #     im = Image.open(imgPath)
        #     ppmPath = f"{imgPath[:imgPath.find('.')]}.ppm"
        #     im.convert("RGB").save(ppmPath, "PPM")
        #     os.rename(imgPath, os.path.join("wrong_formats", image))
        #     imgPath = ppmPath

        # img = cv2.imread(imgPath, cv2.IMREAD_COLOR)
        # img = cv2.resize(img, (IMG_WIDTH, IMG_HEIGHT))
        # output = model.predict(np.array([img]))

        # newPath = os.path.join(id_dir, f"{str(imgs.index(image))}.jpg")
        # im = Image.open(imgPath)
        # im.convert("RGB").save(newPath, "JPEG")

        hierarchy = plants[image[:image.find('.')]]
        if hierarchy == None:
            continue

        slide=prs.slides.add_slide(lyt) #New slide
        title=slide.shapes.title
        title.top = Inches(0)
        title.left = Inches(6)
        imag=slide.shapes.add_picture(imgPath, Inches(0.2), Inches(0.2), width=Inches(2.5), height=Inches(3.333)) #Image
        subtitle=slide.placeholders[1]

        # plantName = data[str(np.argmax(output))]
        plantName = imgPath[:imgPath.find('.')]
        # subtitle.text = plantName #Image classification
        title.text=f"{str(hierarchy.czechName).capitalize()} ({hierarchy.latinName})"
        subtitle.text = f"Třída: {hierarchy.plantClass}\n Řád: {hierarchy.order}\n Čeleď: {hierarchy.family}\n Rod: {hierarchy.genus}"

        # os.remove(newPath)

        # json_object = json.dumps(resultDict, indent=4)
        # with open(valDir, "a") as f:
        #     f.write(json_object)
        prs.save("plants.pptx")

plants = {}
runByList = True
def main():
    if (len(sys.argv) in [2] and sys.argv[1] == "true"):
        eraseFile()
        # getHierarchy("Galega_officinalis")
        # getHierarchy("Zornia gibbosa")
        # getHierarchy("Oncostema_elongata")
        # getHierarchy("Lomatia_ferruginea")

        driver = webdriver.Chrome(ChromeDriverManager().install())
        driver.get(searchUrl)

        consent = driver.find_element(By.ID, 'consentAllButton')
        consent.click()

        if runByList:
            with open(plantsList, "r") as f:
                csvreader = csv.reader(f)
                for i in csvreader:
                    for j in i:
                        plants[j] = getHierarchy(j, driver)
                        print(j)
        else:
            for j in os.listdir(id_dir):
                if j in [".DS_Store", ".gitkeep"]:
                    continue
                plants[j[:j.find(".")]] = getHierarchy(j[:j.find(".")], driver)
                print(j[:j.find(".")])

        json_object = json.dumps(resultDict, indent=4)
        with open(valDir, "a") as f:
            f.write(json_object)

    if not runByList:
        produce_output()
    tm.makeTree()
    #     with open(plantsList, "r") as f:
    #         csvreader = csv.reader(f)
    #         for i in csvreader:
    #             for j in i:
    #                 getHierarchy(j, driver)

    #     json_object = json.dumps(resultDict, indent=4)
    #     with open(valDir, "a") as f:
    #         f.write(json_object)

    # tm.makeTree()

def eraseFile():
    open(valDir, "w").close()

def getHierarchy(plantName, driver):
    sbox = driver.find_element(By.XPATH, "//input[@type='text' and @autofocus='autofocus' and @name='string']")
    sbox.send_keys(plantName)

    submit = driver.find_element(By.XPATH, "//input[@type='submit' and @class='clbutton' and @value=' OK ']")
    submit.click()

    try:
        driver.find_element(By.XPATH, '//*[@id="screen"]/div[3]/div/p/span[1]')
    except:
        # print("Could not find: " + plantName)
        # return
        try:
            firstElem = driver.find_element(By.XPATH, '//*[@id="screen"]/div[5]/div[1]/div/a')
            firstElem.click()
        except NoSuchElementException:
            try:
                secElem = driver.find_element(By.XPATH, '//*[@id="screen"]/div[6]/div[1]/div/a')
                secElem.click()
            except:
                print(f"Error! Could not find {plantName}")
                return

    page_source = driver.page_source
    soup = BS(page_source, features="lxml")

    #Scrape all values
    kingdom = getTypeElement(soup, "říše")
    phylum = getTypeElement(soup, "oddělení")
    plantClass = getTypeElement(soup, "třída")
    order = getTypeElement(soup, "řád")
    family = getTypeElement(soup, "čeleď")
    genus = getTypeElement(soup, "rod")

    try:
        czechName = driver.find_element(By.XPATH, '//*[@id="screen"]/div[3]/div/h1/strong[1]').text
    except:
        czechName = "-"
    try:
        latinName = driver.find_element(By.XPATH, '//*[@id="screen"]/div[3]/div/h1/strong[2]/em').text
    except:
        latinName = "-"
    try:
        latinName = driver.find_element(By.XPATH, '//*[@id="screen"]/div[3]/div/h1/strong/em').text
    except:
        latinName = "-"

    result = plant(kingdom, phylum, plantClass, order, family, genus, czechName, latinName)

    #Create organized output
    familyDict = {family: [f"{genus}NAME:{czechName} \n({latinName})"]}
    orderDict = {order: familyDict}
    classDict = {plantClass: orderDict}
    phylumDict = {phylum: classDict}
    kingdomDict = {kingdom: phylumDict}
    
    if kingdom in iterableKeys(resultDict):
        if phylum in iterableKeys(resultDict[kingdom]):
            if plantClass in iterableKeys(resultDict[kingdom][phylum]):
                if order in iterableKeys(resultDict[kingdom][phylum][plantClass]):
                    if family in iterableKeys(resultDict[kingdom][phylum][plantClass][order]):
                        if genus in resultDict[kingdom][phylum][plantClass][order][family]:
                            return result
                        else:
                            if type(resultDict[kingdom][phylum][plantClass][order][family]) == list:
                                resultDict[kingdom][phylum][plantClass][order][family].append(f"{genus}NAME:{czechName} \n({latinName})")
                            else:
                                resultDict[kingdom][phylum][plantClass][order][family] = [f"{genus}NAME:{czechName} \n({latinName})"]
                    else:
                        resultDict[kingdom][phylum][plantClass][order][family] = [f"{genus}NAME:{czechName} \n({latinName})"]
                else:
                    resultDict[kingdom][phylum][plantClass][order] = familyDict
            else:
                resultDict[kingdom][phylum][plantClass] = orderDict
        else:
            resultDict[kingdom][phylum] = classDict
    else:
        resultDict[kingdom] = phylumDict

    return result

def iterableKeys(diction):
    resList = []
    for key, value in diction.items():
        resList.append(key)
    return resList

def merge_dicts(dict_list):
    result = {}
    for sub_dict in dict_list:
        for key, value in sub_dict.items():
            if isinstance(value, dict):
                result[key] = merge_dicts([result.get(key, {}), value])
            elif key not in result:
                result[key] = value
            elif value != result[key]:
                result[key] = [result[key], value]
    return result

def getElementContent(soup, XPath):
    identifier = soup.find()

def getTypeElement(soup, input):
    identifier = soup.find(string=re.compile(input))
    latin = identifier.find_next("a")
    czech = latin.find_next("b")
    
    if len(latin) > 0:
        latinText = str(latin.contents[0])
        if latinText[:4] == "<em>":
            latinText = latinText[4:-5]
    else:
        latinText = "-"

    czechText = "-"
    if len(czech) > 0:
        czechText = czech.contents[0]
    return f"{czechText} ({latinText})"

if __name__ == "__main__":
    main()